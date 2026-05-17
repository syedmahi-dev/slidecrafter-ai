# engines/python/engine.py
# python-pptx PowerPoint engine

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex string (with or without #) to RGBColor."""
    hex_color = hex_color.lstrip('#')
    r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
    return RGBColor(r, g, b)


def build_presentation(slide_data: dict, output_dir: str = './output', filename: str = 'presentation') -> str:
    """
    Build a .pptx file from structured slide JSON.
    Returns the absolute path to the generated file.
    """
    os.makedirs(output_dir, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.33)   # 16:9 wide
    prs.slide_height = Inches(7.5)

    theme = slide_data.get('theme', {})
    PRIMARY = theme.get('primaryColor', '1E2761').lstrip('#')
    SECONDARY = theme.get('secondaryColor', 'CADCFC').lstrip('#')
    ACCENT = theme.get('accentColor', 'FFFFFF').lstrip('#')
    FONT_TITLE = theme.get('fontTitle', 'Calibri')
    FONT_BODY = theme.get('fontBody', 'Calibri')

    blank_layout = prs.slide_layouts[6]  # blank layout

    for slide_info in slide_data.get('slides', []):
        slide = prs.slides.add_slide(blank_layout)
        slide_type = slide_info.get('type', 'content')
        is_dark = slide_type in ('title', 'closing') or theme.get('darkBackground', False)

        bg_color = PRIMARY if is_dark else 'FFFFFF'
        title_color = ACCENT if is_dark else PRIMARY
        body_color = SECONDARY if is_dark else '333333'

        # Set background
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb(bg_color)

        # Dispatch to slide type builder
        kwargs = dict(
            slide=slide, info=slide_info,
            title_color=title_color, body_color=body_color,
            primary=PRIMARY, secondary=SECONDARY, accent=ACCENT,
            font_title=FONT_TITLE, font_body=FONT_BODY
        )

        if slide_type == 'title':
            _build_title_slide(**kwargs)
        elif slide_type == 'two-column':
            _build_two_column_slide(**kwargs)
        elif slide_type == 'stats':
            _build_stats_slide(**kwargs)
        elif slide_type == 'quote':
            _build_quote_slide(**kwargs)
        elif slide_type == 'closing':
            _build_closing_slide(**kwargs)
        else:
            _build_content_slide(**kwargs)

        # Speaker notes
        if slide_info.get('notes'):
            slide.notes_slide.notes_text_frame.text = slide_info['notes']

    out_path = os.path.abspath(os.path.join(output_dir, f'{filename}.pptx'))
    prs.save(out_path)
    return out_path


# ──────────────────────────────────────────────────────────────
# Slide type builders
# ──────────────────────────────────────────────────────────────

def _add_textbox(slide, text, x, y, w, h, font_name, font_size, bold=False,
                 italic=False, color='000000', align=PP_ALIGN.LEFT, wrap=True):
    from pptx.util import Inches, Pt
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = hex_to_rgb(color)
    return txBox


def _add_rect(slide, x, y, w, h, fill_color, line_color=None):
    from pptx.util import Inches
    from pptx.oxml.ns import qn
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(fill_color)
    if line_color:
        shape.line.color.rgb = hex_to_rgb(line_color)
    else:
        shape.line.fill.background()
    return shape


def _build_title_slide(slide, info, title_color, body_color, font_title, font_body, **_):
    _add_textbox(slide, info.get('title', ''), 0.5, 2.2, 12.3, 2.0,
                 font_title, 44, bold=True, color=title_color, align=PP_ALIGN.CENTER)
    if info.get('subtitle'):
        _add_textbox(slide, info['subtitle'], 0.5, 4.4, 12.3, 0.9,
                     font_body, 20, color=body_color, align=PP_ALIGN.CENTER)


def _build_content_slide(slide, info, title_color, body_color, primary, font_title, font_body, **_):
    _add_textbox(slide, info.get('title', ''), 0.5, 0.3, 12.3, 1.0,
                 font_title, 32, bold=True, color=title_color)
    _add_rect(slide, 0.5, 1.35, 1.5, 0.07, primary)

    if info.get('bullets'):
        from pptx.util import Inches, Pt
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12.3), Inches(5.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(info['bullets']):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f'• {bullet}'
            p.font.size = Pt(16)
            p.font.name = font_body
            p.font.color.rgb = hex_to_rgb(body_color)


def _build_two_column_slide(slide, info, title_color, body_color, primary, secondary, font_title, font_body, **_):
    _add_textbox(slide, info.get('title', ''), 0.5, 0.3, 12.3, 0.9,
                 font_title, 28, bold=True, color=title_color)
    _add_rect(slide, 0.4, 1.4, 5.8, 5.5, secondary)
    _add_textbox(slide, info.get('leftContent', ''), 0.6, 1.6, 5.5, 5.0,
                 font_body, 15, color='222222')
    _add_rect(slide, 6.6, 1.4, 5.8, 5.5, 'F5F5F5', line_color='E0E0E0')
    _add_textbox(slide, info.get('rightContent', ''), 6.8, 1.6, 5.5, 5.0,
                 font_body, 15, color='222222')


def _build_stats_slide(slide, info, title_color, body_color, primary, secondary, accent, font_title, font_body, **_):
    _add_textbox(slide, info.get('title', ''), 0.5, 0.3, 12.3, 0.9,
                 font_title, 28, bold=True, color=title_color)
    stats = info.get('stats', [])[:4]
    count = len(stats)
    if count == 0:
        return
    box_w = 12.33 / count
    for i, stat in enumerate(stats):
        x = 0.5 + i * box_w
        fill = primary if i % 2 == 0 else secondary
        text_color = accent if i % 2 == 0 else primary
        label_color = 'FFFFFF' if i % 2 == 0 else '333333'
        _add_rect(slide, x, 1.6, box_w - 0.2, 4.5, fill)
        _add_textbox(slide, stat.get('value', ''), x, 2.0, box_w - 0.2, 1.8,
                     font_title, 52, bold=True, color=text_color, align=PP_ALIGN.CENTER)
        _add_textbox(slide, stat.get('label', ''), x, 4.1, box_w - 0.2, 0.9,
                     font_body, 14, color=label_color, align=PP_ALIGN.CENTER)


def _build_quote_slide(slide, info, title_color, body_color, font_title, font_body, **_):
    _add_textbox(slide, '\u201C', 0.3, 0.5, 2.0, 1.5,
                 font_title, 100, color='CCCCCC')
    _add_textbox(slide, info.get('quote', ''), 0.8, 1.5, 11.8, 4.0,
                 font_title, 24, italic=True, color=title_color, align=PP_ALIGN.CENTER)
    if info.get('quoteAuthor'):
        _add_textbox(slide, f"— {info['quoteAuthor']}", 0.8, 5.8, 11.8, 0.6,
                     font_body, 14, color=body_color, align=PP_ALIGN.CENTER)


def _build_closing_slide(slide, info, title_color, body_color, font_title, font_body, **_):
    _add_textbox(slide, info.get('title', ''), 0.5, 2.3, 12.3, 1.8,
                 font_title, 40, bold=True, color=title_color, align=PP_ALIGN.CENTER)
    if info.get('subtitle'):
        _add_textbox(slide, info['subtitle'], 0.5, 4.3, 12.3, 0.9,
                     font_body, 18, color=body_color, align=PP_ALIGN.CENTER)
